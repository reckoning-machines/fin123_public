from __future__ import annotations

from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUT_DIR = Path("docs/presentations")
PNG_DIR = OUT_DIR / "fin123_decision_infrastructure_v1_png"
PPTX_PATH = OUT_DIR / "fin123_decision_infrastructure_v1.pptx"
PDF_PATH = OUT_DIR / "fin123_decision_infrastructure_v1.pdf"
README_PATH = OUT_DIR / "README.md"

BG = "0B0F14"
TEXT = "F5F7FA"
SECONDARY = "AAB2BF"
ACCENT = "58A6FF"
PANEL = "111820"
BORDER = "2A3948"
RED = "D86161"

SLIDE_W_IN = 13.333333
SLIDE_H_IN = 7.5
PNG_W = 1920
PNG_H = 1080


def rgb(hex_value: str) -> RGBColor:
    return RGBColor(int(hex_value[0:2], 16), int(hex_value[2:4], 16), int(hex_value[4:6], 16))


def pin(x: float) -> int:
    return Inches(x)


def safe_font(size: int, bold: bool = False):
    return {"name": "Arial", "size": size, "bold": bold}


class SlidePainter:
    def __init__(self, index: int, title: str):
        self.prs = None
        self.slide = None
        self.index = index
        self.title = title
        self.img = Image.new("RGB", (PNG_W, PNG_H), f"#{BG}")
        self.draw = ImageDraw.Draw(self.img)
        self.title_font = self._font(54, bold=True)
        self.body_font = self._font(28)
        self.small_font = self._font(21)
        self.node_font = self._font(27, bold=True)
        self.mono_font = self._font(20)

    def bind(self, prs: Presentation, slide):
        self.prs = prs
        self.slide = slide

    @staticmethod
    def _font(size: int, bold: bool = False):
        candidates = [
            "/System/Library/Fonts/Supplemental/Arial Bold.ttf" if bold else "/System/Library/Fonts/Supplemental/Arial.ttf",
            "/Library/Fonts/Arial Bold.ttf" if bold else "/Library/Fonts/Arial.ttf",
            "/System/Library/Fonts/Helvetica.ttc",
        ]
        for path in candidates:
            try:
                return ImageFont.truetype(path, size)
            except Exception:
                pass
        return ImageFont.load_default()

    def sx(self, x: float) -> int:
        return round(x / SLIDE_W_IN * PNG_W)

    def sy(self, y: float) -> int:
        return round(y / SLIDE_H_IN * PNG_H)

    def rect(self, x: float, y: float, w: float, h: float, fill: str = PANEL, line: str = BORDER):
        shape = self.slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, pin(x), pin(y), pin(w), pin(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill)
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(1)
        self.draw.rectangle([self.sx(x), self.sy(y), self.sx(x + w), self.sy(y + h)], fill=f"#{fill}", outline=f"#{line}", width=2)
        return shape

    def text(self, text: str, x: float, y: float, w: float, h: float, size: int = 24, color: str = TEXT, bold: bool = False, align=PP_ALIGN.LEFT):
        box = self.slide.shapes.add_textbox(pin(x), pin(y), pin(w), pin(h))
        tf = box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.name = "Arial"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)

        font = self._font(round(size * 1.35), bold=bold)
        px, py = self.sx(x), self.sy(y)
        if align == PP_ALIGN.CENTER:
            bbox = self.draw.textbbox((0, 0), text, font=font)
            px = self.sx(x) + max(0, (self.sx(w) - (bbox[2] - bbox[0])) // 2)
        self.draw.text((px, py), text, font=font, fill=f"#{color}")
        return box

    def title_block(self, sentence: str | None = None):
        self.text(self.title, 0.62, 0.42, 9.2, 0.62, 34, TEXT, True)
        self.line(0.62, 1.17, 12.72, 1.17, BORDER, arrow=False)
        if sentence:
            self.text(sentence, 0.68, 6.53, 11.4, 0.36, 18, SECONDARY, False)
        self.text("$fin123", 0.62, 7.08, 1.4, 0.2, 9, SECONDARY, True)
        self.text(f"{self.index:02d}", 12.35, 7.08, 0.55, 0.2, 9, SECONDARY, False, PP_ALIGN.RIGHT)

    def node(self, label: str, x: float, y: float, w: float = 2.0, h: float = 0.62, fill: str = PANEL, color: str = TEXT):
        self.rect(x, y, w, h, fill)
        self.text(label, x + 0.08, y + 0.18, w - 0.16, h - 0.1, 17, color, True, PP_ALIGN.CENTER)

    def placeholder(self, label: str, x: float, y: float, w: float, h: float):
        self.rect(x, y, w, h, "0E141B", BORDER)
        self.text("SCREENSHOT PLACEHOLDER", x + 0.22, y + 0.22, w - 0.44, 0.24, 11, ACCENT, True)
        self.text(label, x + 0.22, y + h / 2 - 0.2, w - 0.44, 0.4, 18, SECONDARY, False, PP_ALIGN.CENTER)

    def line(self, x1: float, y1: float, x2: float, y2: float, color: str = BORDER, arrow: bool = True):
        conn = self.slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, pin(x1), pin(y1), pin(x2), pin(y2))
        conn.line.color.rgb = rgb(color)
        conn.line.width = Pt(1.4)
        self.draw.line([self.sx(x1), self.sy(y1), self.sx(x2), self.sy(y2)], fill=f"#{color}", width=3)
        if arrow:
            # Keep PPTX connectors simple for compatibility; draw arrowheads only in PNG review renders.
            ax, ay = self.sx(x2), self.sy(y2)
            if abs(x2 - x1) > abs(y2 - y1):
                self.draw.polygon([(ax, ay), (ax - 12, ay - 7), (ax - 12, ay + 7)], fill=f"#{color}")
            else:
                self.draw.polygon([(ax, ay), (ax - 7, ay - 12), (ax + 7, ay - 12)], fill=f"#{color}")


def add_blank(prs: Presentation, index: int, title: str) -> SlidePainter:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = rgb(BG)
    p = SlidePainter(index, title)
    p.bind(prs, slide)
    return p


def chain(p: SlidePainter, labels: list[str], x: float, y: float, w: float, gap: float, accent_index: int | None = None):
    for i, label in enumerate(labels):
        p.node(label, x + i * (w + gap), y, w, 0.65, "13263A" if i == accent_index else PANEL, ACCENT if i == accent_index else TEXT)
        if i < len(labels) - 1:
            p.line(x + i * (w + gap) + w, y + 0.32, x + (i + 1) * (w + gap), y + 0.32)


def vertical(p: SlidePainter, labels: list[str], x: float, y: float, w: float, gap: float, accent: set[str] | None = None):
    accent = accent or set()
    for i, label in enumerate(labels):
        p.node(label, x, y + i * gap, w, 0.58, "13263A" if label in accent else PANEL, ACCENT if label in accent else TEXT)
        if i < len(labels) - 1:
            p.line(x + w / 2, y + i * gap + 0.58, x + w / 2, y + (i + 1) * gap)


def build() -> list[SlidePainter]:
    prs = Presentation()
    prs.slide_width = pin(SLIDE_W_IN)
    prs.slide_height = pin(SLIDE_H_IN)
    slides: list[SlidePainter] = []

    p = add_blank(prs, 1, "$fin123")
    p.text("$fin123", 0.66, 0.58, 4.2, 0.62, 38, TEXT, True)
    p.text("Decision Infrastructure for Investment Research", 0.68, 1.34, 8.8, 0.4, 23, SECONDARY, True)
    p.text("Investment decisions became production workflows.", 0.68, 6.12, 8.4, 0.36, 20, TEXT, True)
    p.text("The industry built infrastructure around data.", 0.82, 2.18, 3.25, 0.3, 15, SECONDARY, False, PP_ALIGN.CENTER)
    p.text("The industry built infrastructure around execution.", 4.88, 2.18, 3.25, 0.3, 15, SECONDARY, False, PP_ALIGN.CENTER)
    p.text("$fin123 brings infrastructure to decisions.", 8.9, 2.18, 3.35, 0.3, 15, SECONDARY, False, PP_ALIGN.CENTER)
    p.rect(0.78, 2.72, 3.35, 2.45)
    p.rect(4.84, 2.72, 3.35, 2.45)
    p.rect(8.9, 2.72, 3.35, 2.45, "13263A")
    vertical(p, ["Bloomberg", "FactSet", "Visible Alpha"], 1.22, 3.05, 2.45, 0.58)
    vertical(p, ["OMS", "EMS", "Risk", "Compliance"], 5.28, 2.93, 2.45, 0.52)
    p.node("$fin123", 9.48, 3.55, 2.15, 0.75, "13263A", ACCENT)
    p.text("+", 4.36, 3.75, 0.25, 0.3, 25, SECONDARY, True)
    p.text("+", 8.42, 3.75, 0.25, 0.3, 25, SECONDARY, True)
    slides.append(p)

    p = add_blank(prs, 2, "Data Infrastructure")
    p.title_block("Data workflows became structured, centralized, and governed.")
    p.node("Market Data", 1.05, 2.2, 2.2)
    p.line(3.25, 2.51, 4.05, 2.51)
    p.node("Vendors", 4.05, 2.2, 2.0)
    p.line(6.05, 2.51, 6.85, 2.51)
    p.node("Warehouses", 6.85, 2.2, 2.2)
    p.line(9.05, 2.51, 9.85, 2.51)
    p.node("Research", 9.85, 2.2, 2.0, fill="13263A", color=ACCENT)
    p.text("Bloomberg     FactSet     Visible Alpha     Snowflake     Databases", 1.55, 4.15, 10.2, 0.38, 19, SECONDARY, False, PP_ALIGN.CENTER)
    slides.append(p)

    p = add_blank(prs, 3, "Execution Infrastructure")
    p.title_block("Execution workflows became routed, monitored, and controlled.")
    chain(p, ["OMS", "EMS", "Risk", "Compliance"], 1.55, 3.0, 2.1, 1.0, None)
    p.node("Order Lifecycle", 4.58, 1.86, 4.15, 0.66, "13263A", ACCENT)
    slides.append(p)

    p = add_blank(prs, 4, "Research Still Lives In Files")
    p.title_block("The decision moved. The system of record did not.")
    chain(p, ["Excel", "Word", "PDF", "Email", "Chat"], 0.9, 2.25, 1.65, 0.55)
    p.rect(4.15, 4.05, 5.25, 1.2, "151820", ACCENT)
    p.text("Forecast Changed", 4.4, 4.3, 4.75, 0.34, 25, TEXT, True, PP_ALIGN.CENTER)
    p.text("Why?", 4.4, 4.83, 4.75, 0.28, 19, ACCENT, True, PP_ALIGN.CENTER)
    slides.append(p)

    p = add_blank(prs, 5, "Investment Decisions Are Workflows")
    p.title_block("This is the production workflow behind every investment decision.")
    chain(p, ["Observation", "Research", "Discussion", "Forecast", "Decision", "Position", "Outcome"], 0.55, 3.2, 1.38, 0.27, 4)
    slides.append(p)

    p = add_blank(prs, 6, "One Workflow. Five Stakeholders")
    p.title_block("Each stakeholder touches the same decision workflow.")
    p.node("Decision Model", 5.0, 3.0, 3.2, 0.8, "13263A", ACCENT)
    for label, x, y in [("Analyst", 1.15, 1.95), ("PM", 10.1, 1.95), ("Head of Research", 0.9, 5.0), ("Trading Desk", 9.75, 5.0), ("Compliance", 5.15, 5.85)]:
        p.line(6.6, 3.4, x + 0.95, y + 0.31, arrow=False)
        p.node(label, x, y, 1.9, 0.62)
    slides.append(p)

    p = add_blank(prs, 7, "Analyst")
    p.title_block("Analysts build governed research workflows.")
    p.placeholder("Future fin123 model screenshot", 6.8, 1.8, 5.25, 3.85)
    p.node("Decision Model", 1.35, 3.1, 2.75, 0.76, "13263A", ACCENT)
    for label, x, y in [("Data", 0.95, 2.05), ("AI", 4.2, 2.05), ("Versions", 0.95, 4.78), ("Scenarios", 4.2, 4.78), ("Runs", 2.58, 5.55)]:
        p.line(2.72, 3.48, x + 0.78, y + 0.31, arrow=False)
        p.node(label, x, y, 1.55, 0.58)
    slides.append(p)

    p = add_blank(prs, 8, "Portfolio Manager")
    p.title_block("PMs need provenance when estimates move.")
    p.node("Estimate", 5.25, 3.05, 2.75, 0.76, "13263A", ACCENT)
    for label, x, y in [("Data", 1.1, 2.0), ("Evidence", 3.1, 1.7), ("Memory", 5.35, 1.5), ("AI", 7.75, 1.7), ("Methodology", 9.55, 2.0)]:
        p.line(x + 0.8, y + 0.58, 6.63, 3.05)
        p.node(label, x, y, 1.65, 0.58)
    p.placeholder("Future provenance screenshot", 2.0, 4.85, 9.3, 1.25)
    slides.append(p)

    p = add_blank(prs, 9, "Head of Research")
    p.title_block("Methodology becomes reusable infrastructure.")
    vertical(p, ["House View", "Coverage Universe", "100 Companies", "100 Results"], 4.65, 1.9, 4.05, 1.02, {"House View", "100 Results"})
    slides.append(p)

    p = add_blank(prs, 10, "Trading Desk")
    p.title_block("Real-time market knowledge should enter the decision process.")
    p.text("Current", 0.9, 1.82, 2.0, 0.25, 13, SECONDARY, True)
    chain(p, ["Trader", "Chat", "Phone Call", "Lost"], 0.9, 2.45, 1.7, 0.72, None)
    p.text("Future", 0.9, 4.23, 2.0, 0.25, 13, SECONDARY, True)
    chain(p, ["Trader", "YAP", "Approved Memory", "Forecast", "Decision"], 0.9, 4.86, 1.76, 0.46, 4)
    slides.append(p)

    p = add_blank(prs, 11, "Market Knowledge Should Compound")
    p.title_block("Research compounds. Data compounds. Market knowledge should too.")
    vertical(p, ["Observation", "Discussion", "Review", "Approved Memory", "Future Decisions"], 4.5, 1.65, 4.2, 0.86, {"Approved Memory", "Future Decisions"})
    slides.append(p)

    p = add_blank(prs, 12, "Compliance")
    p.title_block("Compliance needs the ability to reconstruct the decision.")
    vertical(p, ["Data", "Evidence", "Decision", "Replay", "Audit"], 4.95, 1.85, 3.45, 0.88, {"Decision", "Audit"})
    slides.append(p)

    p = add_blank(prs, 13, "Decision Model Lifecycle")
    p.title_block("Governance, memory, replay, and audit live in the lifecycle.")
    vertical(p, ["Decision Model", "Version", "Scenario", "Run", "Result", "Replay", "Audit"], 4.8, 1.45, 3.75, 0.69, {"Decision Model", "Audit"})
    slides.append(p)

    p = add_blank(prs, 14, "Workbook Vision")
    p.title_block("Different surfaces. Same governed substrate.")
    p.node("Workbook", 5.0, 2.05, 3.35, 0.72, "13263A", ACCENT)
    p.line(6.68, 2.77, 4.42, 3.72, arrow=False)
    p.line(6.68, 2.77, 9.02, 3.72, arrow=False)
    p.node("Spreadsheet", 3.2, 3.72, 2.45, 0.66)
    p.node("Document", 7.78, 3.72, 2.45, 0.66)
    slides.append(p)

    p = add_blank(prs, 15, "Why Firms Adopt $fin123")
    p.title_block("$fin123 brings infrastructure to the decision process itself.")
    labels = ["Better Forecasts", "Faster Research", "Institutional Memory", "Replayable Decisions", "Governed AI", "Decision Infrastructure"]
    for i, label in enumerate(labels):
        x = 1.0 + (i % 3) * 4.05
        y = 1.8 + (i // 3) * 1.15
        p.node(label, x, y, 3.2, 0.72, "13263A" if i == 5 else PANEL, ACCENT if i == 5 else TEXT)
    p.text("Data infrastructure. Execution infrastructure. Decision infrastructure.", 1.2, 5.55, 11.0, 0.36, 22, TEXT, True, PP_ALIGN.CENTER)
    slides.append(p)

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    PNG_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(PPTX_PATH)
    for p in slides:
        p.img.save(PNG_DIR / f"slide_{p.index:02d}.png")
    slides[0].img.save(PDF_PATH, save_all=True, append_images=[p.img for p in slides[1:]])
    return slides


def write_readme() -> None:
    README_PATH.write_text(
        """# fin123 Decision Infrastructure v1

Output: `fin123_decision_infrastructure_v1.pptx`

Theme: dark institutional hedge fund deck, Bloomberg meets Apple. Built with `python-pptx` using compatibility-safe rectangles, lines, and text boxes. No animations, SmartArt, video, or custom fonts.

Slide structure:

1. `$fin123` category creation: data infrastructure + execution infrastructure + decision infrastructure.
2. `Data Infrastructure`: data layer workflow.
3. `Execution Infrastructure`: order lifecycle workflow.
4. `Research Still Lives In Files`: files fail to explain forecast changes.
5. `Investment Decisions Are Workflows`: observation to outcome motif.
6. `One Workflow. Five Stakeholders`: hub-and-spoke decision model.
7. `Analyst`: build workflow with screenshot placeholder.
8. `Portfolio Manager`: estimate provenance with screenshot placeholder.
9. `Head of Research`: house view scales across coverage.
10. `Trading Desk`: current lost knowledge vs future approved memory.
11. `Market Knowledge Should Compound`: observation to future decisions.
12. `Compliance`: reconstructable decision chain.
13. `Decision Model Lifecycle`: ontology from model to audit.
14. `Workbook Vision`: spreadsheet and document on governed substrate.
15. `Why Firms Adopt $fin123`: adoption outcomes and final thesis.

Review exports:

PNG previews are in `fin123_decision_infrastructure_v1_png/`.

PDF export:

`fin123_decision_infrastructure_v1.pdf` is generated from the same slide layout as the PNG previews for quick review.
""",
        encoding="utf-8",
    )


if __name__ == "__main__":
    built = build()
    write_readme()
    print(f"Wrote {PPTX_PATH}")
    print(f"Wrote {PDF_PATH}")
    print(f"Wrote {README_PATH}")
    print(f"Wrote {len(built)} PNG previews to {PNG_DIR}")
